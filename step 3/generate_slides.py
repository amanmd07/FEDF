import sys
import os

# Ensure python-pptx is installed. If not, try installing it.
try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Attempting to install...")
    import subprocess
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("Successfully installed python-pptx!")
    except Exception as e:
        print(f"Failed to install python-pptx: {e}")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Theme colors
    BG_DARK = RGBColor(9, 14, 26)
    CARD_BG = RGBColor(21, 32, 54)
    TEXT_MAIN = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    PRIMARY = RGBColor(99, 102, 241)     # Indigo
    SECONDARY = RGBColor(14, 165, 233)   # Light Blue
    ACCENT = RGBColor(16, 185, 129)      # Emerald
    BORDER_COLOR = RGBColor(38, 48, 71)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    def apply_background(slide):
        # Add full-screen rectangle for background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background() # No border
        return bg

    def add_header_footer(slide, title_text, slide_num, total_slides=5):
        # Slide Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        
        # Header Accent Line
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = SECONDARY
        accent_line.line.fill.background()

        # Footer Text
        footerBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.3))
        ftf = footerBox.text_frame
        ftf.word_wrap = True
        ftf.margin_left = ftf.margin_top = ftf.margin_right = ftf.margin_bottom = 0
        fp = ftf.paragraphs[0]
        fp.text = f"Hospital Bed Management Platform  |  Slide {slide_num} of {total_slides}"
        fp.font.name = 'Arial'
        fp.font.size = Pt(10)
        fp.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover Page)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)

    # Large decorative accent block on the left
    left_banner = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_banner.fill.solid()
    left_banner.fill.fore_color.rgb = PRIMARY
    left_banner.line.fill.background()

    # Title & Subtitle in single text frame to avoid overlapping
    mainBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(3))
    tf1 = mainBox.text_frame
    tf1.word_wrap = True
    
    p_badge = tf1.paragraphs[0]
    p_badge.text = "INSEM-END PROJECT PRESENTATION"
    p_badge.font.name = 'Arial'
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = SECONDARY
    p_badge.space_after = Pt(20)

    p_title = tf1.add_paragraph()
    p_title.text = "Real-Time Hospital Bed\nManagement Dashboard"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    p_title.space_after = Pt(12)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Front-End & Back-End Architectural Overview (Slides 10 - 12)"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = TEXT_MUTED

    # Metadata at bottom
    metaBox = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(11), Inches(1.5))
    mtf = metaBox.text_frame
    mtf.word_wrap = True
    
    mp1 = mtf.paragraphs[0]
    mp1.text = "Submitted By: MD Aman (2520030190) & V. Sree Pragnya (2520030238)"
    mp1.font.name = 'Arial'
    mp1.font.size = Pt(13)
    mp1.font.bold = True
    mp1.font.color.rgb = TEXT_MAIN
    mp1.space_after = Pt(4)

    mp2 = mtf.add_paragraph()
    mp2.text = "Course: Front End Development Frameworks  |  Under Guidance: Dr. Ramya Krishana Dhulipalla"
    mp2.font.name = 'Arial'
    mp2.font.size = Pt(12)
    mp2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Front-End Architecture & Core Usability (Page 10)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_header_footer(slide2, "Front-End Architecture & Core Usability", 2)

    # Left Column: Key Features
    leftBox = slide2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8))
    ltf = leftBox.text_frame
    ltf.word_wrap = True
    
    lp_desc = ltf.paragraphs[0]
    lp_desc.text = "A responsive React-based interface engineered for hospital operations, optimizing speed and reducing workflow errors."
    lp_desc.font.name = 'Arial'
    lp_desc.font.size = Pt(15)
    lp_desc.font.color.rgb = TEXT_MUTED
    lp_desc.space_after = Pt(20)

    features = [
        ("Component & Navigation Consistency", "A unified navigation bar houses the project brand, menu links, and logout options to maintain system state."),
        ("Medical-Theme Styling System", "Clean CSS forms, tables, and buttons help operators quickly distinguish key actions (admission vs. discharge)."),
        ("Multi-Device Responsiveness", "Fully responsive layouts ensure seamless utilization across desktop PCs, laptop stations, and mobile tablets.")
    ]

    for title, desc in features:
        p_title = ltf.add_paragraph()
        p_title.text = f"•  {title}"
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = SECONDARY
        
        p_desc = ltf.add_paragraph()
        p_desc.text = f"    {desc}"
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_after = Pt(12)

    # Right Column: Dashboard Role Cards
    card_width = Inches(5.2)
    card_height = Inches(1.3)
    card_left = Inches(7.3)
    
    roles = [
        ("Receptionist Dashboard", "Gives complete control over patient admission, bed allocation, and discharge forms.", PRIMARY),
        ("Nurse Dashboard", "Provides real-time monitors to track bed vacancies and active patient coordinates.", SECONDARY),
        ("Doctor Dashboard", "Offers read-only access to patient health logs and bed occupation records.", ACCENT)
    ]

    for i, (title, desc, color) in enumerate(roles):
        card_top = Inches(1.7 + i * 1.6)
        
        # Draw Card Background
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)
        
        # Indicator strip on the left of each card
        strip = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_left, card_top, Inches(0.12), card_height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()

        # Text Frame
        ctfBox = slide2.shapes.add_textbox(card_left + Inches(0.3), card_top + Inches(0.15), card_width - Inches(0.4), card_height - Inches(0.3))
        ctf = ctfBox.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = 0
        
        cp_title = ctf.paragraphs[0]
        cp_title.text = title
        cp_title.font.name = 'Arial'
        cp_title.font.size = Pt(14)
        cp_title.font.bold = True
        cp_title.font.color.rgb = TEXT_MAIN
        cp_title.space_after = Pt(2)
        
        cp_desc = ctf.add_paragraph()
        cp_desc.text = desc
        cp_desc.font.name = 'Arial'
        cp_desc.font.size = Pt(11)
        cp_desc.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 3: Front-End Components & Key Design Features (Page 11)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_header_footer(slide3, "Front-End Components & Interface Design", 3)

    # Subtitle or description
    descBox = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
    dtf = descBox.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = "Key UI screens and functional assets mapped to user requirements:"
    dp.font.name = 'Arial'
    dp.font.size = Pt(14)
    dp.font.color.rgb = TEXT_MUTED

    # Table
    rows, cols = 6, 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(6.1)
    table.columns[2].width = Inches(2.8)

    headers = ["Page / Component", "Description & Operational Purpose", "Key Attributes"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = SECONDARY

    comp_data = [
        ("Landing Page", "Serves as the application's entry screen and introduces the core platform metrics.", "Brand introduction"),
        ("Login Portal", "Gives secure, credential-gated access with custom selections for user roles.", "Secure JWT gating"),
        ("Receptionist Dashboard", "Supports forms for admitting patients, allocating beds, and discharging patients.", "Operational focus"),
        ("Nurse / Doctor Dashboard", "Allows medical staff to monitor patient listings and active bed allocations.", "Monitoring & Logs"),
        ("UI Elements (Forms/Buttons)", "Reusable data input forms and action buttons (Submit, Assign, Discharge).", "Consistent workflow")
    ]

    for row_idx, (comp, desc, attr) in enumerate(comp_data, start=1):
        # Component column
        cell_comp = table.cell(row_idx, 0)
        cell_comp.text = comp
        cell_comp.fill.solid()
        cell_comp.fill.fore_color.rgb = BG_DARK
        p_comp = cell_comp.text_frame.paragraphs[0]
        p_comp.font.name = 'Arial'
        p_comp.font.size = Pt(12)
        p_comp.font.bold = True
        p_comp.font.color.rgb = TEXT_MAIN

        # Description column
        cell_desc = table.cell(row_idx, 1)
        cell_desc.text = desc
        cell_desc.fill.solid()
        cell_desc.fill.fore_color.rgb = BG_DARK
        p_desc = cell_desc.text_frame.paragraphs[0]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

        # Attributes column
        cell_attr = table.cell(row_idx, 2)
        cell_attr.text = attr
        cell_attr.fill.solid()
        cell_attr.fill.fore_color.rgb = BG_DARK
        p_attr = cell_attr.text_frame.paragraphs[0]
        p_attr.font.name = 'Arial'
        p_attr.font.size = Pt(11)
        p_attr.font.bold = True
        p_attr.font.color.rgb = ACCENT

    # -------------------------------------------------------------
    # SLIDE 4: Back-End Design & Middleware Security (Page 11-12)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_header_footer(slide4, "Back-End Design & Middleware Security", 4)

    # Left Column: Backend Architecture Notes
    leftBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8))
    ltf4 = leftBox4.text_frame
    ltf4.word_wrap = True
    
    lp_desc4 = ltf4.paragraphs[0]
    lp_desc4.text = "Built on Node.js and Express.js, the backend services handle all application business logic, authorization checks, and data files."
    lp_desc4.font.name = 'Arial'
    lp_desc4.font.size = Pt(15)
    lp_desc4.font.color.rgb = TEXT_MUTED
    lp_desc4.space_after = Pt(20)

    be_features = [
        ("Modular API Architecture", "Different server modules handle distinct responsibilities: auth controller, patient registry, and bed states."),
        ("Role-Based Middleware Gate", "Intercepts calls, validates JWT signatures, and prevents unauthorized access to protected methods."),
        ("Atomic State Changes", "Discharging a patient automatically resets bed status parameters and frees slots for allocation.")
    ]

    for title, desc in be_features:
        p_title = ltf4.add_paragraph()
        p_title.text = f"•  {title}"
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = SECONDARY
        
        p_desc = ltf4.add_paragraph()
        p_desc.text = f"    {desc}"
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_after = Pt(12)

    # Right Column: CSS Diagram Representation
    # We will build a visual flow-chart using pptx shapes
    diag_left = Inches(7.3)
    diag_width = Inches(5.2)
    diag_height = Inches(4.2)
    
    # Diagram container background
    container = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, diag_left, Inches(1.8), diag_width, diag_height)
    container.fill.solid()
    container.fill.fore_color.rgb = CARD_BG
    container.line.color.rgb = BORDER_COLOR
    container.line.width = Pt(1)

    # Diagram Title
    dtBox = slide4.shapes.add_textbox(diag_left + Inches(0.2), Inches(2.0), diag_width - Inches(0.4), Inches(0.4))
    dtf = dtBox.text_frame
    dtp = dtf.paragraphs[0]
    dtp.text = "API Middleware Request Life Cycle"
    dtp.alignment = PP_ALIGN.CENTER
    dtp.font.name = 'Arial'
    dtp.font.size = Pt(14)
    dtp.font.bold = True
    dtp.font.color.rgb = TEXT_MAIN

    # Flow elements
    # 1. Client Request
    step1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diag_left + Inches(0.4), Inches(2.8), Inches(1.2), Inches(1.0))
    step1.fill.solid()
    step1.fill.fore_color.rgb = BG_DARK
    step1.line.color.rgb = PRIMARY
    step1.text = "Client\nRequest"
    step1.text_frame.paragraphs[0].font.size = Pt(11)
    step1.text_frame.paragraphs[0].font.bold = True

    # Arrow 1
    arr1 = slide4.shapes.add_textbox(diag_left + Inches(1.65), Inches(3.1), Inches(0.3), Inches(0.4))
    arr1.text_frame.text = "→"
    arr1.text_frame.paragraphs[0].font.size = Pt(20)
    arr1.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # 2. Middleware Access Check
    step2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diag_left + Inches(1.95), Inches(2.8), Inches(1.3), Inches(1.0))
    step2.fill.solid()
    step2.fill.fore_color.rgb = PRIMARY
    step2.line.color.rgb = BORDER_COLOR
    step2.text = "Auth\nMiddleware"
    step2.text_frame.paragraphs[0].font.size = Pt(11)
    step2.text_frame.paragraphs[0].font.bold = True

    # Arrow 2
    arr2 = slide4.shapes.add_textbox(diag_left + Inches(3.25), Inches(3.1), Inches(0.3), Inches(0.4))
    arr2.text_frame.text = "→"
    arr2.text_frame.paragraphs[0].font.size = Pt(20)
    arr2.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # 3. Target API Controller
    step3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, diag_left + Inches(3.55), Inches(2.8), Inches(1.25), Inches(1.0))
    step3.fill.solid()
    step3.fill.fore_color.rgb = BG_DARK
    step3.line.color.rgb = ACCENT
    step3.text = "API\nRoute"
    step3.text_frame.paragraphs[0].font.size = Pt(11)
    step3.text_frame.paragraphs[0].font.bold = True

    # Text below diagram
    infoBox = slide4.shapes.add_textbox(diag_left + Inches(0.4), Inches(4.3), diag_width - Inches(0.8), Inches(1.4))
    itf = infoBox.text_frame
    itf.word_wrap = True
    ip = itf.paragraphs[0]
    ip.text = "Gatekeeping Security Enforced:"
    ip.font.name = 'Arial'
    ip.font.size = Pt(12)
    ip.font.bold = True
    ip.font.color.rgb = TEXT_MAIN
    ip.space_after = Pt(4)
    
    ip2 = itf.add_paragraph()
    ip2.text = "Middleware interceptors (middleware/auth.js) inspect the HTTP headers, extracting and parsing the JWT. If the token signature is invalid, or if the parsed user role lacks access to the requested path (e.g. nurse discharging a patient), the request is aborted."
    ip2.font.name = 'Arial'
    ip2.font.size = Pt(10.5)
    ip2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 5: Core Back-End Modules & APIs (Page 12)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_header_footer(slide5, "Back-End Modules & Server Routing", 5)

    # Subtitle or description
    descBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
    dtf5 = descBox5.text_frame
    dtf5.word_wrap = True
    dp5 = dtf5.paragraphs[0]
    dp5.text = "Express modules and database interface mappings:"
    dp5.font.name = 'Arial'
    dp5.font.size = Pt(14)
    dp5.font.color.rgb = TEXT_MUTED

    # Table
    rows5, cols5 = 6, 3
    left5, top5, width5, height5 = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
    table_shape5 = slide5.shapes.add_table(rows5, cols5, left5, top5, width5, height5)
    table5 = table_shape5.table
    
    table5.columns[0].width = Inches(2.8)
    table5.columns[1].width = Inches(6.1)
    table5.columns[2].width = Inches(2.8)

    # Headers
    headers5 = ["API Module", "Functional Responsibility", "File Path / Route"]
    for i, h in enumerate(headers5):
        cell = table5.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Arial'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = SECONDARY

    be_data = [
        ("Server Entry", "Initializes Express server, binds port, mounts middleware, and routes API requests.", "server/index.js"),
        ("Auth Controller", "Accepts login requests, queries credentials, and signs secure tokens.", "server/routes/auth.js"),
        ("Role Middleware", "Enforces access control policies by checking verified roles on active routes.", "server/middleware/auth.js"),
        ("Patient Registry", "Exposes API endpoints for patient admission, records searches, and discharge commands.", "server/routes/patients.js"),
        ("Bed Registry", "Maintains database records for bed locations, capacities, and assignment states.", "server/routes/beds.js")
    ]

    for row_idx, (comp, desc, attr) in enumerate(be_data, start=1):
        # Module column
        cell_comp = table5.cell(row_idx, 0)
        cell_comp.text = comp
        cell_comp.fill.solid()
        cell_comp.fill.fore_color.rgb = BG_DARK
        p_comp = cell_comp.text_frame.paragraphs[0]
        p_comp.font.name = 'Arial'
        p_comp.font.size = Pt(12)
        p_comp.font.bold = True
        p_comp.font.color.rgb = TEXT_MAIN

        # Responsibility column
        cell_desc = table5.cell(row_idx, 1)
        cell_desc.text = desc
        cell_desc.fill.solid()
        cell_desc.fill.fore_color.rgb = BG_DARK
        p_desc = cell_desc.text_frame.paragraphs[0]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED

        # Route column
        cell_attr = table5.cell(row_idx, 2)
        cell_attr.text = attr
        cell_attr.fill.solid()
        cell_attr.fill.fore_color.rgb = BG_DARK
        p_attr = cell_attr.text_frame.paragraphs[0]
        p_attr.font.name = 'Arial'
        p_attr.font.size = Pt(11)
        p_attr.font.bold = True
        p_attr.font.color.rgb = ACCENT

    # Save presentation to Desktop
    desktop_path = r"C:\Users\DELL\OneDrive\Desktop"
    out_file = os.path.join(desktop_path, "Hospital_Dashboard_Presentation.pptx")
    prs.save(out_file)
    print(f"Successfully generated PowerPoint presentation at: {out_file}")

if __name__ == "__main__":
    create_presentation()
